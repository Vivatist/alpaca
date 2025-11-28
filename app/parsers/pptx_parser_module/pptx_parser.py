#!/usr/bin/env python3
"""
PowerPoint Presentation Parser для RAG системы ALPACA

Парсер PowerPoint презентаций (.pptx) без OCR.

Pipeline:
    .pptx → Unstructured → Markdown + YAML metadata

Возможности:
- Извлечение текста из слайдов с сохранением структуры
- Поддержка русскоязычных и английских презентаций
- Извлечение метаданных (автор, количество слайдов)
- Генерация YAML header с метаданными
- Экспорт в Markdown формат для RAG индексации
- БЕЗ OCR для изображений
"""

import os
import sys
from pathlib import Path
from typing import Dict, Optional

# Добавляем путь к базовому парсеру
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from base_parser import BaseParser

try:
    from unstructured.partition.pptx import partition_pptx  # type: ignore
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False


class PptxParser(BaseParser):
    """
    Парсер PowerPoint презентаций без OCR
    
    Использует pipeline:
    1. Unstructured - основной парсер текста из слайдов
    2. python-pptx - извлечение метаданных
    3. YAML header - метаданные для RAG
    """
    
    def __init__(self):
        """Инициализация парсера"""
        super().__init__("pptx-parser")
        
        # Проверка зависимостей
        self._check_dependencies()
        
    def _check_dependencies(self):
        """Проверка наличия необходимых библиотек"""
        missing = []
        
        if not UNSTRUCTURED_AVAILABLE:
            missing.append("unstructured[pptx]")
        if not PYTHON_PPTX_AVAILABLE:
            missing.append("python-pptx")
            
        if missing:
            self.logger.warning(f"Missing dependencies | packages={', '.join(missing)}")
    
    def parse(self, file_path: str, file_hash: Optional[str] = None) -> Dict:
        """
        Парсинг PowerPoint презентации в Markdown с метаданными
        
        Args:
            file_path: Путь к .pptx файлу
            file_hash: Хэш файла от file-watcher (опционально)
            
        Returns:
            Dict с ключами:
                - markdown: str - Текст в Markdown формате
                - metadata: Dict - Метаданные документа
                - yaml_header: str - YAML header для документа
                - success: bool - Статус парсинга
                - error: Optional[str] - Сообщение об ошибке
        """
        result = {
            'markdown': '',
            'metadata': {},
            'yaml_header': '',
            'success': False,
            'error': None
        }
        
        try:
            if not os.path.exists(file_path):
                result['error'] = f"File not found: {file_path}"
                return result
            
            self.logger.info(f"Parsing PowerPoint presentation | file={file_path}")
            
            # 1. Добавляем ОБЩИЕ метаданные (в базовом классе)
            common_metadata = self._add_common_metadata(file_path, file_hash)
            
            # 2. Извлечение СПЕЦИФИЧНЫХ метаданных PowerPoint через python-pptx
            specific_metadata = self._extract_pptx_specific_metadata(file_path)
            
            # 3. Основной парсинг через Unstructured
            markdown_content = self._parse_with_unstructured(file_path)
            
            result['markdown'] = markdown_content
            
            # 4. Объединяем все метаданные для результата
            result['metadata'] = {**common_metadata, **specific_metadata}
            
            # 5. Генерация YAML header с разделением общих и специфичных метаданных
            yaml_header = self._generate_yaml_header(
                common_metadata,
                specific_metadata,
                'pptx'
            )
            result['yaml_header'] = yaml_header
            
            result['success'] = True
            
            self.logger.info(f"PowerPoint presentation parsed successfully | file={file_path} slides={specific_metadata.get('slides', 0)}")
            
        except Exception as e:
            result['error'] = str(e)
            result['success'] = False
            self.logger.error(f"Error parsing PowerPoint presentation | file={file_path} error={type(e).__name__}: {e}")
        
        return result
    
    def _extract_pptx_specific_metadata(self, file_path: str) -> Dict:
        """
        Извлечение СПЕЦИФИЧНЫХ для PowerPoint метаданных
        
        Общие метаданные (file_name, file_path, file_size, etc.) добавляются в базовом классе.
        Здесь добавляем только специфичные для PPTX данные.
        
        Args:
            file_path: Путь к файлу
            
        Returns:
            Dict со специфичными метаданными
        """
        specific_metadata = {
            'author': '',
            'title': '',
            'subject': '',
            'slides': 0,
            'ocr_enabled': False
        }
        
        if not PYTHON_PPTX_AVAILABLE:
            return specific_metadata
        
        try:
            prs = Presentation(file_path)
            core_props = prs.core_properties
            
            specific_metadata['author'] = core_props.author or ''
            specific_metadata['title'] = core_props.title or ''
            specific_metadata['subject'] = core_props.subject or ''
            specific_metadata['slides'] = len(prs.slides)
            
            self.logger.debug(f"PPTX-specific metadata | author={specific_metadata['author']} slides={specific_metadata['slides']} title={specific_metadata['title']}")
            
        except Exception as e:
            self.logger.warning(f"PPTX-specific metadata extraction failed | file={file_path} error={type(e).__name__}: {e}")

        return specific_metadata
    
    def _parse_with_unstructured(self, file_path: str) -> str:
        """
        Парсинг презентации через Unstructured
        
        Args:
            file_path: Путь к файлу
            
        Returns:
            Markdown текст
        """
        if not UNSTRUCTURED_AVAILABLE:
            self.logger.warning("Unstructured not available, using fallback parser")
            return self._fallback_parse(file_path)
        
        try:
            self.logger.info(f"Parsing with Unstructured | file={file_path}")
            
            # Используем partition_pptx с поддержкой русского и английского языков
            elements = partition_pptx(
                filename=file_path,
                infer_table_structure=True,
                languages=["rus", "eng"]  # Поддержка русского и английского
            )
            
            if not elements:
                self.logger.warning(f"No elements extracted | file={file_path}")
                return self._fallback_parse(file_path)
            
            # Группируем элементы по слайдам
            markdown_parts = []
            current_slide = None
            slide_content = []
            
            for element in elements:
                # Пытаемся определить номер слайда из метаданных
                metadata = getattr(element, 'metadata', None)
                slide_number = None
                
                if metadata:
                    # Unstructured может хранить номер страницы/слайда
                    slide_number = getattr(metadata, 'page_number', None)
                
                # Если номер слайда изменился, сохраняем предыдущий слайд
                if slide_number is not None and slide_number != current_slide:
                    if slide_content:
                        markdown_parts.append("\n\n".join(slide_content))
                    
                    current_slide = slide_number
                    slide_content = [f"## Слайд {slide_number}"]
                
                # Добавляем текст элемента
                element_text = str(element).strip()
                if element_text:
                    slide_content.append(element_text)
            
            # Добавляем последний слайд
            if slide_content:
                markdown_parts.append("\n\n".join(slide_content))
            
            markdown = "\n\n---\n\n".join(markdown_parts)
            
            self.logger.info(f"Unstructured parsing complete | slides={len(markdown_parts)} text_length={len(markdown)}")
            
            return markdown
            
        except Exception as e:
            self.logger.error(f"Unstructured parsing failed | file={file_path} error={type(e).__name__}: {e}")
            return self._fallback_parse(file_path)
    
    def _fallback_parse(self, file_path: str) -> str:
        """
        Резервный парсер через python-pptx
        
        Args:
            file_path: Путь к файлу
            
        Returns:
            Простой текст презентации
        """
        if not PYTHON_PPTX_AVAILABLE:
            return "ERROR: No parser available (install unstructured[pptx] or python-pptx)"
        
        try:
            self.logger.info(f"Using fallback parser (python-pptx) | file={file_path}")
            
            prs = Presentation(file_path)
            markdown_parts = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_parts = [f"## Слайд {slide_idx}"]
                
                # Извлекаем текст из всех shapes на слайде
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Пытаемся определить заголовки
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text:
                                    # Если большой шрифт или жирный - вероятно заголовок
                                    if paragraph.level == 0:
                                        slide_parts.append(f"### {para_text}")
                                    else:
                                        slide_parts.append(para_text)
                        else:
                            slide_parts.append(text)
                    
                    # Извлекаем текст из таблиц
                    if hasattr(shape, "table"):
                        table_md = self._table_to_markdown(shape.table)
                        if table_md:
                            slide_parts.append(table_md)
                
                if len(slide_parts) > 1:  # Есть контент кроме заголовка слайда
                    markdown_parts.append("\n\n".join(slide_parts))
            
            markdown = "\n\n---\n\n".join(markdown_parts)
            
            self.logger.info(f"Fallback parsing complete | slides={len(markdown_parts)} text_length={len(markdown)}")
            
            return markdown
            
        except Exception as e:
            self.logger.error(f"Fallback parser failed | file={file_path} error={type(e).__name__}: {e}")
            return f"ERROR: Failed to parse presentation: {str(e)}"
    
    def _table_to_markdown(self, table) -> str:
        """
        Конвертация таблицы PowerPoint в Markdown
        
        Args:
            table: Таблица из python-pptx
            
        Returns:
            Markdown представление таблицы
        """
        try:
            if not table.rows:
                return ""
            
            lines = []
            
            # Заголовок таблицы (первая строка)
            header_cells = [cell.text.strip() for cell in table.rows[0].cells]
            lines.append("| " + " | ".join(header_cells) + " |")
            lines.append("|" + "|".join(["---"] * len(header_cells)) + "|")
            
            # Остальные строки
            for row in table.rows[1:]:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append("| " + " | ".join(cells) + " |")
            
            return "\n".join(lines)
            
        except Exception as e:
            self.logger.warning(f"Failed to convert table to markdown | error={type(e).__name__}: {e}")
            return ""


if __name__ == "__main__":
    # Тестовый запуск
    import argparse
    
    parser = argparse.ArgumentParser(description="PPTX Parser Test")
    parser.add_argument("file", help="Path to PPTX file")
    parser.add_argument("--output", "-o", help="Output markdown file")
    
    args = parser.parse_args()
    
    pptx_parser = PptxParser()
    result = pptx_parser.parse(args.file)
    
    if result['success']:
        print("✓ Parsing successful")
        print(f"Slides: {result['metadata'].get('slides', 0)}")
        print(f"Text length: {len(result['markdown'])} chars")
        
        if args.output:
            pptx_parser.save_to_markdown_file(result, args.output)
            print(f"✓ Saved to {args.output}")
        else:
            print("\n" + "="*60)
            print(result['yaml_header'])
            print(result['markdown'][:500])
            print("...")
    else:
        print(f"✗ Parsing failed: {result['error']}")
        sys.exit(1)
