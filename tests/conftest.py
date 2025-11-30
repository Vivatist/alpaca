"""
Pytest fixtures для тестирования воркера
"""
import os
import tempfile
import pytest
import psycopg2
import logging
from datetime import datetime
from pathlib import Path
from typing import Generator

from alpaca.infrastructure.database.postgres import PostgresFileRepository
from settings import settings


@pytest.fixture(scope="session", autouse=True)
def setup_test_logging():
    """Настройка логирования для тестов - изолированно от основного приложения"""
    # Настраиваем логирование для тестов
    root_logger = logging.getLogger()
    
    # Очищаем все существующие handlers
    root_logger.handlers.clear()
    
    # Создаём новый handler для тестов
    handler = logging.StreamHandler()
    handler.setLevel(logging.ERROR)  # В тестах показываем только ошибки
    formatter = logging.Formatter('%(asctime)s | %(levelname)-8s | %(name)s - %(message)s')
    handler.setFormatter(formatter)
    root_logger.addHandler(handler)
    root_logger.setLevel(logging.ERROR)
    
    yield
    
    # После тестов не очищаем handlers - main.py уже настроил своё логирование
    # перед запуском тестов и продолжит использовать его


@pytest.fixture
def test_db() -> Generator[PostgresFileRepository, None, None]:
    """Создание тестовой БД"""
    db = PostgresFileRepository(settings.DATABASE_URL)
    yield db
    # Cleanup после тестов
    with db.get_connection() as conn:
        with conn.cursor() as cur:
            # Очистка тестовых данных
            cur.execute("DELETE FROM chunks WHERE metadata->>'file_path' LIKE '/tmp/test_%'")
            cur.execute("DELETE FROM files WHERE path LIKE '/tmp/test_%'")
        conn.commit()


@pytest.fixture
def temp_test_file() -> Generator[tuple[str, str], None, None]:
    """Создание временного тестового файла"""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', prefix='test_', delete=False) as f:
        test_content = "Это тестовый файл.\n\nОн содержит несколько параграфов.\n\nДля проверки работы чанкования."
        f.write(test_content)
        temp_path = f.name
    
    yield temp_path, test_content
    
    # Cleanup
    if os.path.exists(temp_path):
        os.unlink(temp_path)


@pytest.fixture
def temp_docx_file() -> Generator[str, None, None]:
    """Создание временного DOCX файла для тестирования парсинга"""
    from docx import Document
    
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.docx', prefix='test_', delete=False) as f:
        doc = Document()
        doc.add_paragraph("Заголовок тестового документа")
        doc.add_paragraph("Это первый параграф с тестовым содержимым.")
        doc.add_paragraph("Это второй параграф.")
        doc.add_paragraph("А это третий параграф с дополнительным текстом для проверки парсинга.")
        doc.save(f.name)
        temp_path = f.name
    
    yield temp_path
    
    # Cleanup
    if os.path.exists(temp_path):
        os.unlink(temp_path)


@pytest.fixture
def temp_pptx_file() -> Generator[str, None, None]:
    """Создание временного PPTX файла для тестов парсера презентаций"""
    from pptx import Presentation
    from pptx.util import Inches

    with tempfile.NamedTemporaryFile(mode='wb', suffix='.pptx', prefix='test_', delete=False) as f:
        prs = Presentation()

        # Первый слайд с заголовком и списком
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Тестовая презентация"
        body_shape = slide.shapes.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.text = "Первый пункт"
        paragraph = text_frame.add_paragraph()
        paragraph.text = "Вложенный пункт"
        paragraph.level = 1

        # Второй слайд с таблицей
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = "Таблица"
        table = slide2.shapes.add_table(2, 2, Inches(0.5), Inches(1.5), Inches(9), Inches(2.5)).table
        table.cell(0, 0).text = "Колонка A"
        table.cell(0, 1).text = "Колонка B"
        table.cell(1, 0).text = "Значение 1"
        table.cell(1, 1).text = "Значение 2"

        prs.save(f.name)
        temp_path = f.name

    yield temp_path

    if os.path.exists(temp_path):
        os.unlink(temp_path)


@pytest.fixture
def temp_xlsx_file() -> Generator[str, None, None]:
    """Создание временного XLSX файла для тестов парсера Excel"""
    from openpyxl import Workbook

    with tempfile.NamedTemporaryFile(mode='wb', suffix='.xlsx', prefix='test_', delete=False) as f:
        wb = Workbook()

        sheet = wb.active
        sheet.title = "Сводка"
        sheet.append(["Раздел", "Значение", "Комментарий"])
        sheet.append(["Доход", 1250000.75, "рост 12%"])
        sheet.append(["Расход", 842300.10, "снижение 3%"])
        sheet.append(["План", None, "утвержден 2025-11-01"])

        sheet2 = wb.create_sheet("Детализация")
        sheet2['A1'] = "Код"
        sheet2['B1'] = "Описание"
        sheet2['C1'] = "Дата"
        sheet2.append([101, "Бурение", datetime(2025, 11, 1)])
        sheet2.append([205, "Логистика", datetime(2025, 11, 5)])

        wb.save(f.name)
        temp_path = f.name

    yield temp_path

    if os.path.exists(temp_path):
        os.unlink(temp_path)


@pytest.fixture
def temp_xls_file() -> Generator[str, None, None]:
    """Создание временного XLS файла для тестов наследия"""
    import xlwt

    fd, temp_path = tempfile.mkstemp(suffix='.xls', prefix='test_legacy_')
    os.close(fd)

    try:
        wb = xlwt.Workbook()
        sheet = wb.add_sheet("План")

        date_style = xlwt.easyxf(num_format_str='YYYY-MM-DD')
        datetime_style = xlwt.easyxf(num_format_str='YYYY-MM-DD HH:MM:SS')

        sheet.write(0, 0, "Этап")
        sheet.write(0, 1, "Дата")
        sheet.write(0, 2, "Сумма")
        sheet.write(1, 0, "Старт")
        sheet.write(1, 1, datetime(2025, 1, 5), date_style)
        sheet.write(1, 2, 750000.0)
        sheet.write(2, 0, "Подготовка")
        sheet.write(2, 1, datetime(2025, 1, 7, 14, 30), datetime_style)
        sheet.write(2, 2, 125000.5)

        sheet2 = wb.add_sheet("Раскладка")
        sheet2.write(0, 0, "Категория")
        sheet2.write(0, 1, "Количество")
        sheet2.write(1, 0, "Блюда")
        sheet2.write(1, 1, 12)

        wb.save(temp_path)
        yield temp_path
    finally:
        if os.path.exists(temp_path):
            os.unlink(temp_path)


@pytest.fixture
def mock_file_info():
    """Создание мока информации о файле"""
    def _create_mock(file_path: str, file_hash: str, status: str = 'added'):
        return {
            'path': file_path,
            'hash': file_hash,
            'status_sync': status
        }
    return _create_mock


@pytest.fixture
def cleanup_temp_parsed():
    """Очистка временных распарсенных файлов"""
    temp_dir = "/home/alpaca/tmp_md"
    yield
    # Cleanup после тестов
    if os.path.exists(temp_dir):
        for file in Path(temp_dir).rglob("test_*"):
            if file.is_file():
                file.unlink()
